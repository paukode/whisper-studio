"""create_docx / create_pptx / create_xlsx / create_pdf.

These exercise the REAL builders (python-docx, python-pptx, openpyxl,
reportlab) rather than mocking them — the whole point of these tools is that
they produce genuine, round-trippable files, so a mocked subprocess/library
call would not catch a broken invocation.

The tools never ask the user where to save. With no workspace connected the
destination comes from destination_path (the user's own words) or defaults
to ~/Documents/<filename>; the finished bytes are STAGED in the app sandbox
and the [WS_APPROVAL] (action=save_to_path) payload carries staged_path —
approving just moves the ready file into place. With a workspace connected
it's action=create_document keyed by the workspace-relative path, same
human-in-the-loop gate ws_create_file goes through. The actual writes are
exercised separately against server.approval.bootstrap._do_create_document /
_do_save_to_path.
"""

import asyncio
import base64
import io
import json
import os
import re
import zipfile
from unittest.mock import patch

import pytest

from server.approval.bootstrap import _do_create_document, _do_save_to_path
from server.documents.executors import (
    _exec_create_docx,
    _exec_create_pdf,
    _exec_create_pptx,
    _exec_create_xlsx,
)
from server.workspace.paths import is_staged_path


def _run(coro):
    return asyncio.run(coro)


def _parse_prompt(output: str) -> dict:
    assert output.startswith("[WS_WORKSPACE_PROMPT]"), output
    return json.loads(output[len("[WS_WORKSPACE_PROMPT]") :])


def _parse_approval(output: str, expected_action: str = "create_document") -> dict:
    assert output.startswith("[WS_APPROVAL]"), output
    payload = json.loads(output[len("[WS_APPROVAL]") :])
    assert payload["action"] == expected_action
    return payload


def _decoded_bytes(payload: dict) -> bytes:
    if "content_b64" in payload:
        return base64.b64decode(payload["content_b64"])
    with open(payload["staged_path"], "rb") as f:
        return f.read()


# ── No workspace connected: default to ~/Documents, stage, never ask ──


def test_create_docx_no_workspace_defaults_to_documents_and_stages(monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: None)
    out = _exec_create_docx({"path": "report.docx", "html_content": "<p>hi</p>"}, "", {})
    payload = _parse_approval(out, expected_action="save_to_path")
    assert payload["path"] == os.path.expanduser("~/Documents/report.docx")
    assert is_staged_path(payload["staged_path"])
    assert _decoded_bytes(payload)[:2] == b"PK", "staged file must be real OOXML (zip)"


def test_create_docx_appends_missing_extension(monkeypatch):
    # Models sometimes pass a bare name like 'document' — the deliverable
    # must still open in Word.
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: None)
    out = _exec_create_docx({"path": "document", "html_content": "<p>hi</p>"}, "", {})
    payload = _parse_approval(out, expected_action="save_to_path")
    assert payload["path"] == os.path.expanduser("~/Documents/document.docx")
    assert payload["filename"] == "document.docx"


def test_create_pptx_no_workspace_defaults_and_stages(monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: None)
    out = _exec_create_pptx(
        {"path": "deck.pptx", "slides": [{"title": "T", "bullets": ["a"]}]}, "", {}
    )
    payload = _parse_approval(out, expected_action="save_to_path")
    assert payload["path"] == os.path.expanduser("~/Documents/deck.pptx")
    assert is_staged_path(payload["staged_path"])


def test_create_xlsx_no_workspace_defaults_and_stages(monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: None)
    out = _exec_create_xlsx(
        {"path": "data.xlsx", "sheets": [{"name": "Sheet1", "rows": [[1, 2]]}]}, "", {}
    )
    payload = _parse_approval(out, expected_action="save_to_path")
    assert payload["path"] == os.path.expanduser("~/Documents/data.xlsx")
    assert is_staged_path(payload["staged_path"])


def test_create_pdf_no_workspace_defaults_and_stages(monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: None)
    out = _exec_create_pdf({"path": "summary.pdf", "title": "T", "paragraphs": ["hello"]}, "", {})
    payload = _parse_approval(out, expected_action="save_to_path")
    assert payload["path"] == os.path.expanduser("~/Documents/summary.pdf")
    assert _decoded_bytes(payload)[:5] == b"%PDF-"


def test_create_docx_with_destination_path_from_user_words(tmp_path, monkeypatch):
    """When the model resolved a destination from the user's words, the tool
    builds and stages in one call — get_workspace_path() is never consulted,
    matching save_file's own no-workspace-required path."""
    monkeypatch.setattr(
        "server.workspace.state.get_workspace_path",
        lambda: (_ for _ in ()).throw(AssertionError("must not check workspace")),
    )
    dest = str(tmp_path / "report.docx")
    out = _exec_create_docx(
        {"path": "report.docx", "html_content": "<p>hi</p>", "destination_path": dest}, "", {}
    )
    payload = _parse_approval(out, expected_action="save_to_path")
    assert payload["path"] == os.path.realpath(dest)
    assert payload["filename"] == "report.docx"
    assert is_staged_path(payload["staged_path"])
    assert not (tmp_path / "report.docx").exists(), "must not write before approval"


def test_create_docx_rejects_relative_destination_path(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: None)
    out = _exec_create_docx(
        {"path": "report.docx", "html_content": "<p>hi</p>", "destination_path": "relative.docx"},
        "",
        {},
    )
    assert out.startswith("Error"), out


def test_do_save_to_path_moves_staged_document_on_approval(tmp_path):
    from server.workspace.paths import stage_file_bytes

    data = b"fake docx bytes for the direct save flow"
    staged = stage_file_bytes(data, "notes.docx")
    dest = str(tmp_path / "notes.docx")
    outcome = _run(
        _do_save_to_path({"path": dest, "filename": "notes.docx", "staged_path": staged})
    )
    assert outcome.ok, outcome.error
    assert (tmp_path / "notes.docx").read_bytes() == data


def test_create_docx_path_outside_workspace_emits_sentinel(tmp_path, monkeypatch):
    ws = tmp_path / "ws"
    ws.mkdir()
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(ws))
    # The destination is resolved BEFORE any rendering work: an unusable path
    # must come back as the folder-picker sentinel without building anything.
    with patch("server.documents.executors.render_html") as mock_render:
        out = _exec_create_docx({"path": "../outside.docx", "html_content": "<p>hi</p>"}, "", {})
    mock_render.assert_not_called()
    payload = _parse_prompt(out)
    assert payload["reason"] == "outside_workspace"


# ── Real conversions: build succeeds, pauses for approval, never touches disk ──


def test_create_docx_produces_real_ooxml_and_pauses_for_approval(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = _exec_create_docx(
        {
            "path": "report.docx",
            "html_content": "<h1>Title</h1><p>Body text with <b>bold</b>.</p>",
        },
        "",
        {},
    )
    payload = _parse_approval(out)
    assert payload["path"] == "report.docx"
    assert payload["format"] == "docx"
    assert not (tmp_path / "report.docx").exists(), "must not write before approval"

    data = _decoded_bytes(payload)
    assert payload["size"] == len(data) > 0
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        names = z.namelist()
        assert "word/document.xml" in names
        assert "[Content_Types].xml" in names
        doc_xml = z.read("word/document.xml").decode("utf-8")
        assert "Body text" in doc_xml


def test_create_docx_applies_standard_layout(tmp_path, monkeypatch):
    """The standard: A4 page, 1-inch margins, Calibri 11pt body."""
    from docx import Document
    from docx.shared import Inches, Pt

    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = _exec_create_docx(
        {"path": "std.docx", "html_content": "<h1>Heading</h1><p>Body paragraph.</p>"}, "", {}
    )
    doc = Document(io.BytesIO(_decoded_bytes(_parse_approval(out))))

    section = doc.sections[0]
    # A4 is 210x297mm; python-docx rounds to EMU, so compare in millimetres.
    assert round(section.page_width.mm) == 210
    assert round(section.page_height.mm) == 297
    assert section.left_margin == Inches(1) and section.top_margin == Inches(1)
    assert doc.styles["Normal"].font.name == "Calibri"
    assert doc.styles["Normal"].font.size == Pt(11)

    body = next(p for p in doc.paragraphs if p.text == "Body paragraph.")
    assert body.runs[0].font.name == "Calibri"
    assert body.runs[0].font.size == Pt(11)
    assert doc.styles["Heading 1"].font.size == Pt(16)


def test_create_docx_renders_the_supported_markup(tmp_path, monkeypatch):
    """The advertised vocabulary must actually land as native Word
    constructs, not as flattened text."""
    from docx import Document

    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    html = (
        "<h1>Title</h1><h2>Sub</h2>"
        "<p>Plain <b>bold</b> <i>italic</i> <u>under</u> <code>mono</code>.</p>"
        "<ul><li>One</li><li>Two<ul><li>Nested</li></ul></li></ul>"
        "<ol><li>Step</li></ol>"
        "<table><tr><th>H</th></tr><tr><td>C</td></tr></table>"
    )
    out = _exec_create_docx({"path": "rich.docx", "html_content": html}, "", {})
    doc = Document(io.BytesIO(_decoded_bytes(_parse_approval(out))))

    styles = [p.style.name for p in doc.paragraphs if p.text.strip()]
    assert styles[:2] == ["Heading 1", "Heading 2"]
    assert "List Bullet" in styles and "List Bullet 2" in styles and "List Number" in styles

    runs = {r.text: r for r in doc.paragraphs[2].runs}
    assert runs["bold"].bold and runs["italic"].italic and runs["under"].underline
    assert runs["mono"].font.name == "Consolas"

    table = doc.tables[0]
    assert (len(table.rows), len(table.columns)) == (2, 1)
    assert table.rows[0].cells[0].paragraphs[0].runs[0].bold, "th must be bold"


@pytest.mark.parametrize(
    "cell_html, expected",
    [
        ('<td bgcolor="#E63946">x</td>', "E63946"),
        ('<td bgcolor="E63946">x</td>', "E63946"),
        ('<td style="background-color: rgb(0,128,0)">x</td>', "008000"),
        ('<td style="background: yellow">x</td>', "FFFF00"),
        ('<td style="background-color:#abc">x</td>', "AABBCC"),
    ],
)
def test_create_docx_honors_cell_background_colors(cell_html, expected, tmp_path, monkeypatch):
    """The regression that exposed the whole gap: the previous textutil
    converter dropped cell shading silently, so a user asking for a colored
    table got an uncolored one with no error."""
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = _exec_create_docx(
        {"path": "c.docx", "html_content": f"<table><tr>{cell_html}</tr></table>"}, "", {}
    )
    with zipfile.ZipFile(io.BytesIO(_decoded_bytes(_parse_approval(out)))) as z:
        doc_xml = z.read("word/document.xml").decode("utf-8")
    assert re.findall(r'<w:shd[^>]*w:fill="([0-9A-F]{6})"', doc_xml) == [expected]


def test_create_docx_survives_malformed_markup(tmp_path, monkeypatch):
    """Model-written HTML is not guaranteed well-formed. An unclosed tag
    should cost formatting on the tail, never the document."""
    from docx import Document

    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = _exec_create_docx(
        {"path": "m.docx", "html_content": "<p>open <b>bold<p>next</p></div>"}, "", {}
    )
    doc = Document(io.BytesIO(_decoded_bytes(_parse_approval(out))))
    text = " ".join(p.text for p in doc.paragraphs)
    assert "open" in text and "bold" in text and "next" in text


def test_create_docx_rejects_empty_html(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = _exec_create_docx({"path": "empty.docx", "html_content": "   "}, "", {})
    assert out.startswith("Error"), out
    assert not (tmp_path / "empty.docx").exists()


def test_create_pptx_round_trips_with_python_pptx(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    slides = [
        {"title": "Intro", "bullets": ["First point", "Second point"]},
        {"title": "Conclusion", "bullets": ["Wrap up"]},
    ]
    out = _exec_create_pptx({"path": "deck.pptx", "slides": slides}, "", {})
    payload = _parse_approval(out)
    assert payload["format"] == "pptx"
    assert not (tmp_path / "deck.pptx").exists(), "must not write before approval"
    data = _decoded_bytes(payload)
    assert len(data) > 0

    from pptx import Presentation

    from server.documents.standards import PPTX_SLIDE_HEIGHT_EMU, PPTX_SLIDE_WIDTH_EMU

    prs = Presentation(io.BytesIO(data))
    assert prs.slide_width == PPTX_SLIDE_WIDTH_EMU, "slides must be standard 16:9"
    assert prs.slide_height == PPTX_SLIDE_HEIGHT_EMU
    # Placeholders must span the widescreen slide, not the old 4:3 layout box.
    title_shape = prs.slides[0].shapes.title
    assert title_shape.width == prs.slide_width - 2 * title_shape.left
    assert len(prs.slides) == 2
    titles = [s.shapes.title.text for s in prs.slides]
    assert titles == ["Intro", "Conclusion"]
    first_slide_bullets = [p.text for p in prs.slides[0].placeholders[1].text_frame.paragraphs]
    assert first_slide_bullets == ["First point", "Second point"]


def test_create_xlsx_round_trips_with_openpyxl(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    sheets = [
        {"name": "Sales", "rows": [["Item", "Qty"], ["Widget", 3]]},
        {"name": "Notes", "rows": [["free text"]]},
    ]
    out = _exec_create_xlsx({"path": "data.xlsx", "sheets": sheets}, "", {})
    payload = _parse_approval(out)
    assert payload["format"] == "xlsx"
    assert not (tmp_path / "data.xlsx").exists(), "must not write before approval"
    data = _decoded_bytes(payload)
    assert len(data) > 0

    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data))
    assert wb.sheetnames == ["Sales", "Notes"]
    ws = wb["Sales"]
    assert ws["A1"].value == "Item"
    assert ws["B1"].value == "Qty"
    assert ws["A2"].value == "Widget"
    assert ws["B2"].value == 3
    assert wb["Notes"]["A1"].value == "free text"
    # Standard polish: bold frozen header row, fitted column widths.
    assert ws["A1"].font.bold is True
    assert ws.freeze_panes == "A2"
    assert (ws.column_dimensions["A"].width or 0) >= 10


def test_create_pdf_produces_valid_pdf_bytes(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = _exec_create_pdf(
        {
            "path": "summary.pdf",
            "title": "Quarterly Summary",
            "paragraphs": ["Revenue is up.", "Costs < projections & on track."],
        },
        "",
        {},
    )
    payload = _parse_approval(out)
    assert payload["format"] == "pdf"
    assert not (tmp_path / "summary.pdf").exists(), "must not write before approval"
    data = _decoded_bytes(payload)
    assert len(data) > 0
    assert data[:5] == b"%PDF-"

    # Standard layout: A4 pages (595.28 x 841.89 pt).
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(data)
    width, height = pdf[0].get_size()
    assert abs(width - 595.28) < 1 and abs(height - 841.89) < 1


def test_create_pdf_escapes_special_characters_without_dropping_content(tmp_path, monkeypatch):
    """A plain paragraph containing '<', '>', '&' must survive verbatim —
    reportlab's Paragraph treats unescaped input as its own mini-XML markup
    and will silently swallow or mangle text that looks like a tag."""
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    tricky = "<unknown>tag</unknown> AT&T reported 5 < 10 items"
    out = _exec_create_pdf({"path": "tricky.pdf", "title": "", "paragraphs": [tricky]}, "", {})
    payload = _parse_approval(out)
    data = _decoded_bytes(payload)

    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(data)
    text = pdf[0].get_textpage().get_text_range()
    assert "<unknown>tag</unknown>" in text
    assert "AT&T reported 5 < 10 items" in text


def test_create_pdf_requires_title_or_paragraphs(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = _exec_create_pdf({"path": "empty.pdf", "title": "", "paragraphs": []}, "", {})
    assert out.startswith("Error"), out
    assert not (tmp_path / "empty.pdf").exists()


@pytest.mark.parametrize(
    "executor,tool_input",
    [
        (_exec_create_pptx, {"path": "bad.pptx", "slides": []}),
        (_exec_create_xlsx, {"path": "bad.xlsx", "sheets": []}),
    ],
)
def test_empty_collection_inputs_are_rejected(tmp_path, monkeypatch, executor, tool_input):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = executor(tool_input, "", {})
    assert out.startswith("Error"), out


# ── Approval click actually performs the write ─────────────────────────────


def test_do_create_document_writes_file_on_approval(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: str(tmp_path))
    data = b"fake docx bytes for approval-write test"
    payload = {
        "path": "approved.docx",
        "content_b64": base64.b64encode(data).decode("ascii"),
        "format": "docx",
        "size": len(data),
    }
    outcome = _run(_do_create_document(payload))
    assert outcome.ok, outcome.error
    full = tmp_path / "approved.docx"
    assert full.is_file()
    assert full.read_bytes() == data
    assert "Created approved.docx" in (outcome.output or "")


def test_do_create_document_no_workspace_fails(monkeypatch):
    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: None)
    outcome = _run(_do_create_document({"path": "x.docx", "content_b64": ""}))
    assert not outcome.ok
    assert "No workspace connected" in (outcome.error or "")


def test_do_create_document_rejects_path_outside_workspace(tmp_path, monkeypatch):
    ws = tmp_path / "ws"
    ws.mkdir()
    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: str(ws))
    outcome = _run(
        _do_create_document(
            {"path": "../outside.docx", "content_b64": base64.b64encode(b"x").decode("ascii")}
        )
    )
    assert not outcome.ok
    assert not (tmp_path / "outside.docx").exists()


def test_save_and_create_outcomes_carry_the_revision_hint(tmp_path, monkeypatch):
    """Both write paths (staged save with no workspace, direct create in a
    workspace) must tell the model to send the next revision to the same
    path instead of minting a new filename."""
    from server.workspace.paths import stage_file_bytes
    from server.workspace.state import REVISION_HINT

    staged = stage_file_bytes(b"bytes", "notes.docx")
    dest = str(tmp_path / "notes.docx")
    outcome = _run(
        _do_save_to_path({"path": dest, "filename": "notes.docx", "staged_path": staged})
    )
    assert outcome.ok, outcome.error
    assert REVISION_HINT in outcome.output

    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: str(tmp_path))
    payload = {
        "path": "ws.docx",
        "content_b64": base64.b64encode(b"bytes").decode("ascii"),
        "format": "docx",
        "size": 5,
    }
    outcome = _run(_do_create_document(payload))
    assert outcome.ok, outcome.error
    assert REVISION_HINT in outcome.output
