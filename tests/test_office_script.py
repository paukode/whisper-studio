"""office_script + inspect_document — full-control Office authoring/editing.

These drive the REAL libraries and the REAL sandbox runner rather than
mocking them. The whole reason the tool exists is that the schema-driven
create_* tools cannot express everything users ask for, so a test that
mocked the execution would assert nothing about whether the formatting
actually survives into the file. test_shading_survives_the_round_trip
pins that end to end, against the OOXML itself.

Split of responsibilities under test:
- _exec_office_script validates and packages; it must NEVER run code.
- run_office_script executes approved code in a throwaway dir and returns
  bytes; it must never touch the user's source file or the destination.
- _do_office_script re-resolves the destination at click time and is the
  only thing that writes.
"""

import asyncio
import json
import os
import re
import zipfile

import pytest

from server.documents.approval import _do_office_script
from server.documents.executors import _exec_inspect_document, _exec_office_script
from server.documents.inspect import inspect_document
from server.documents.script import run_office_script, validate_script_request


def _run(coro):
    return asyncio.run(coro)


def _parse_approval(output: str) -> dict:
    assert output.startswith("[WS_APPROVAL]"), output
    payload = json.loads(output[len("[WS_APPROVAL]") :])
    assert payload["action"] == "office_script"
    return payload


def _shaded_fills(data: bytes) -> list[str]:
    xml = zipfile.ZipFile(__import__("io").BytesIO(data)).read("word/document.xml").decode()
    return re.findall(r'<w:shd[^>]*w:fill="([0-9A-Fa-f]{6})"', xml)


# ── Fixtures: stand-in "documents the user already has" ────────────────


@pytest.fixture
def grid_docx(tmp_path):
    """A Word file shaped like the one that motivated this feature: a table
    whose cells are meant to be colored into a picture."""
    from docx import Document

    doc = Document()
    doc.add_heading("Grid", level=1)
    table = doc.add_table(rows=4, cols=4)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "corner"
    path = tmp_path / "grid.docx"
    doc.save(str(path))
    return str(path)


@pytest.fixture
def sales_xlsx(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Region", "Q1", "Q2"])
    ws.append(["North", 10, 12])
    ws.merge_cells("A1:C1")
    ws.freeze_panes = "A2"
    path = tmp_path / "sales.xlsx"
    wb.save(str(path))
    return str(path)


@pytest.fixture
def deck_pptx(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Quarterly review"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return str(path)


# ── inspect_document ───────────────────────────────────────────────────


def test_inspect_docx_reports_addressable_table_coordinates(grid_docx):
    report = inspect_document(grid_docx)
    assert "Tables: 1" in report
    assert "Table 0 — 4 rows x 4 cols" in report
    # The coordinates must be pasteable straight into an office_script call.
    assert "doc.tables[0].rows[r].cells[c]" in report
    assert "Heading 1: Grid" in report
    assert "corner" in report


def test_inspect_xlsx_reports_ranges_merges_and_freeze(sales_xlsx):
    report = inspect_document(sales_xlsx)
    assert "Sheet 'Sales'" in report
    assert "frozen panes at A2" in report
    assert "A1:C1" in report
    assert "'Region'" in report


def test_inspect_pptx_reports_slide_size_and_shapes(deck_pptx):
    report = inspect_document(deck_pptx)
    assert "Slides: 1" in report
    assert "Slide 0" in report
    assert "Quarterly review" in report
    assert "prs.slides[0].shapes[j]" in report


def test_inspect_rejects_legacy_binary_formats(tmp_path):
    legacy = tmp_path / "old.doc"
    legacy.write_bytes(b"\xd0\xcf\x11\xe0")
    assert "not a .docx/.xlsx/.pptx" in inspect_document(str(legacy))


def test_inspect_missing_file_is_an_error_not_a_crash(tmp_path):
    assert inspect_document(str(tmp_path / "nope.docx")).startswith("Error:")


def test_inspect_corrupt_file_is_an_error_not_a_crash(tmp_path):
    broken = tmp_path / "broken.docx"
    broken.write_bytes(b"not a zip at all")
    assert inspect_document(str(broken)).startswith("Error reading")


def test_inspect_executor_resolves_workspace_relative(grid_docx, monkeypatch):
    monkeypatch.setattr(
        "server.workspace.state.get_workspace_path", lambda: os.path.dirname(grid_docx)
    )
    out = _exec_inspect_document({"path": "grid.docx"}, "", {})
    assert "Table 0 — 4 rows x 4 cols" in out


def test_inspect_executor_rejects_relative_path_without_workspace(monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: None)
    out = _exec_inspect_document({"path": "grid.docx"}, "", {})
    assert "must be an absolute path" in out


# ── The executor validates and packages, and never executes ────────────


def test_exec_packages_approval_without_running_anything(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    marker = tmp_path / "side_effect.txt"
    out = _exec_office_script(
        {
            "path": "out.docx",
            "summary": "color the header row",
            "code": f"open({str(marker)!r}, 'w').write('ran')",
        },
        "",
        {},
    )
    payload = _parse_approval(out)
    assert payload["dest_kind"] == "workspace"
    assert payload["path"] == "out.docx"
    assert payload["format"] == "docx"
    assert payload["summary"] == "color the header row"
    assert not marker.exists(), "the executor must not run the code before approval"
    assert not (tmp_path / "out.docx").exists(), "must not write before approval"


def test_exec_no_workspace_defaults_to_documents(monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: None)
    out = _exec_office_script({"path": "out.xlsx", "summary": "s", "code": "pass"}, "", {})
    payload = _parse_approval(out)
    assert payload["dest_kind"] == "absolute"
    assert payload["path"] == os.path.expanduser("~/Documents/out.xlsx")
    assert payload["format"] == "xlsx"


def test_exec_takes_format_from_destination_path_not_bare_path(tmp_path, monkeypatch):
    """`path` is only a suggested filename once destination_path is set, so a
    bare 'out' must not defeat format detection."""
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: None)
    dest = str(tmp_path / "deck.pptx")
    out = _exec_office_script(
        {"path": "out", "summary": "s", "code": "pass", "destination_path": dest}, "", {}
    )
    payload = _parse_approval(out)
    assert payload["format"] == "pptx"
    assert payload["path"] == os.path.realpath(dest)


@pytest.mark.parametrize(
    "tool_input, expected",
    [
        ({"path": "out.docx", "summary": "s", "code": "  "}, "code is required"),
        ({"path": "out.pdf", "summary": "s", "code": "pass"}, "must end in .docx"),
        ({"path": "out", "summary": "s", "code": "pass"}, "must end in .docx"),
    ],
)
def test_exec_rejects_bad_requests_before_asking_the_user(
    tool_input, expected, tmp_path, monkeypatch
):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = _exec_office_script(tool_input, "", {})
    assert out.startswith("Error"), out
    assert expected in out


def test_exec_rejects_unreadable_source_before_asking_the_user(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    out = _exec_office_script(
        {"path": "out.docx", "summary": "s", "code": "pass", "source_path": "missing.docx"},
        "",
        {},
    )
    assert "no such file" in out


def test_exec_rejects_legacy_source_format(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.state.get_workspace_path", lambda: str(tmp_path))
    (tmp_path / "old.doc").write_bytes(b"\xd0\xcf\x11\xe0")
    out = _exec_office_script(
        {"path": "out.docx", "summary": "s", "code": "pass", "source_path": "old.doc"}, "", {}
    )
    assert "must be a .docx/.xlsx/.pptx file" in out


def test_validate_rejects_oversized_code():
    from server.documents.script import MAX_CODE_BYTES

    assert "too large" in validate_script_request("x" * (MAX_CODE_BYTES + 1), "docx")


# ── The runner ─────────────────────────────────────────────────────────


def test_run_creates_a_document_from_scratch():
    code = (
        "from docx import Document\n"
        "doc = Document()\n"
        "doc.add_paragraph('made by a script')\n"
        "print('wrote one paragraph')\n"
        "doc.save(OUTPUT_PATH)\n"
    )
    ok, message, data = run_office_script(code, "docx")
    assert ok, message
    assert message == "wrote one paragraph"
    assert data[:2] == b"PK", "must be real OOXML"


def test_input_path_is_none_when_creating_from_scratch():
    code = (
        "assert INPUT_PATH is None, INPUT_PATH\n"
        "from docx import Document\n"
        "Document().save(OUTPUT_PATH)\n"
    )
    ok, message, _ = run_office_script(code, "docx")
    assert ok, message


def test_shading_survives_the_round_trip(grid_docx):
    """Shading applied by a script must reach the OOXML. create_docx can
    color cells too, but only in a table it is building from scratch; this
    is the path for a document the user already has."""
    code = (
        "from docx import Document\n"
        "from docx.oxml import OxmlElement\n"
        "from docx.oxml.ns import qn\n"
        "doc = Document(INPUT_PATH)\n"
        "for cell in doc.tables[0].rows[0].cells:\n"
        "    shd = OxmlElement('w:shd')\n"
        "    shd.set(qn('w:val'), 'clear')\n"
        "    shd.set(qn('w:fill'), 'E63946')\n"
        "    cell._tc.get_or_add_tcPr().append(shd)\n"
        "doc.save(OUTPUT_PATH)\n"
    )
    ok, message, data = run_office_script(code, "docx", grid_docx)
    assert ok, message
    assert _shaded_fills(data) == ["E63946"] * 4


def test_editing_never_mutates_the_users_source_file(grid_docx):
    """The script gets a COPY. A script that opens and re-saves in place must
    not reach the original, since nobody has approved anything yet."""
    before = open(grid_docx, "rb").read()
    code = (
        "from docx import Document\n"
        "doc = Document(INPUT_PATH)\n"
        "doc.add_paragraph('appended')\n"
        "doc.save(INPUT_PATH)\n"  # deliberately saves over its input
        "doc.save(OUTPUT_PATH)\n"
    )
    ok, message, data = run_office_script(code, "docx", grid_docx)
    assert ok, message
    assert open(grid_docx, "rb").read() == before, "the user's file was modified"
    assert len(data) != len(before) or data != before


def test_run_reports_a_script_error_back_to_the_model():
    ok, message, data = run_office_script("raise ValueError('bad column')", "docx")
    assert not ok
    assert "bad column" in message
    assert data == b""


def test_run_reports_a_script_that_forgot_to_save():
    ok, message, data = run_office_script("print('did nothing')", "docx")
    assert not ok
    assert "OUTPUT_PATH" in message
    assert data == b""


def test_run_rejects_an_empty_output_file():
    ok, message, _ = run_office_script("open(OUTPUT_PATH, 'wb').close()", "docx")
    assert not ok
    assert "empty file" in message


def test_run_edits_a_workbook_with_formatting_create_xlsx_cannot_express(sales_xlsx):
    code = (
        "from openpyxl import load_workbook\n"
        "from openpyxl.styles import PatternFill\n"
        "wb = load_workbook(INPUT_PATH)\n"
        "ws = wb['Sales']\n"
        "ws['B2'].fill = PatternFill('solid', fgColor='FFFF00')\n"
        "ws['B2'].number_format = '#,##0.00'\n"
        "wb.save(OUTPUT_PATH)\n"
    )
    ok, message, data = run_office_script(code, "xlsx", sales_xlsx)
    assert ok, message

    import io

    from openpyxl import load_workbook

    ws = load_workbook(io.BytesIO(data))["Sales"]
    assert ws["B2"].fill.fgColor.rgb.endswith("FFFF00")
    assert ws["B2"].number_format == "#,##0.00"


# ── The approval click is the only thing that writes ───────────────────


def _payload(tmp_path, code, dest_kind="workspace", path="out.docx", source_path=""):
    return {
        "action": "office_script",
        "path": path,
        "dest_kind": dest_kind,
        "format": "docx",
        "code": code,
        "source_path": source_path,
        "summary": "s",
    }


_MAKE_DOCX = "from docx import Document\nDocument().save(OUTPUT_PATH)\n"


def test_approval_writes_into_the_workspace(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: str(tmp_path))
    outcome = _run(_do_office_script(_payload(tmp_path, _MAKE_DOCX)))
    assert outcome.ok, outcome.error
    assert (tmp_path / "out.docx").read_bytes()[:2] == b"PK"
    assert "Wrote out.docx" in outcome.output


def test_approval_writes_to_an_absolute_destination(tmp_path):
    dest = str(tmp_path / "sub" / "out.docx")
    outcome = _run(_do_office_script(_payload(tmp_path, _MAKE_DOCX, "absolute", dest)))
    assert outcome.ok, outcome.error
    assert os.path.isfile(dest), "parent directories must be created"


def test_approval_surfaces_script_output_to_the_model(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: str(tmp_path))
    code = "print('shaded 12 cells')\n" + _MAKE_DOCX
    outcome = _run(_do_office_script(_payload(tmp_path, code)))
    assert outcome.ok, outcome.error
    assert "shaded 12 cells" in outcome.output


def test_approval_writes_nothing_when_the_script_fails(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: str(tmp_path))
    outcome = _run(_do_office_script(_payload(tmp_path, "raise RuntimeError('nope')")))
    assert not outcome.ok
    assert "nope" in outcome.error
    assert not (tmp_path / "out.docx").exists()


def test_approval_refuses_a_path_escaping_the_workspace(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: str(tmp_path))
    outcome = _run(_do_office_script(_payload(tmp_path, _MAKE_DOCX, path="../escaped.docx")))
    assert not outcome.ok
    assert outcome.error == "Invalid path"


def test_approval_refuses_a_relative_absolute_destination(tmp_path):
    outcome = _run(_do_office_script(_payload(tmp_path, _MAKE_DOCX, "absolute", "out.docx")))
    assert not outcome.ok
    assert "absolute path" in outcome.error


def test_approval_refuses_a_sensitive_destination(tmp_path):
    dest = os.path.expanduser("~/.ssh/out.docx")
    outcome = _run(_do_office_script(_payload(tmp_path, _MAKE_DOCX, "absolute", dest)))
    assert not outcome.ok
    assert "sensitive path" in outcome.error


def test_approval_without_a_workspace_is_refused_not_guessed(tmp_path, monkeypatch):
    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: None)
    outcome = _run(_do_office_script(_payload(tmp_path, _MAKE_DOCX)))
    assert not outcome.ok
    assert outcome.error == "No workspace connected"


# ── The tool is actually reachable ─────────────────────────────────────


def test_tools_are_registered_and_routed():
    from server.approval.bootstrap import register_defaults
    from server.approval.registry import get
    from server.documents.tools import DOCUMENT_TOOLS
    from server.executors import EXECUTORS, is_concurrent_safe

    names = {t["name"] for t in DOCUMENT_TOOLS}
    assert {"office_script", "inspect_document"} <= names
    assert "office_script" in EXECUTORS and "inspect_document" in EXECUTORS
    # Reading a document is safe to run alongside other tools; writing is not.
    assert is_concurrent_safe("inspect_document")
    assert not is_concurrent_safe("office_script")

    register_defaults()
    spec = get("office_script")
    assert spec is not None
    # Its own category: approving a document write must not silently
    # authorize running code, nor vice versa.
    assert spec.category == "office-script"
    assert spec.preview == "command"
    assert spec.render_command({"code": "print(1)"}) == "print(1)"


def test_approval_summary_distinguishes_editing_from_building():
    from server.documents.approval import _summary

    assert _summary({"path": "a/b.docx", "summary": "color cells"}) == "Build b.docx — color cells"
    edit = _summary({"path": "b.docx", "summary": "color cells", "source_path": "b.docx"})
    assert edit == "Edit b.docx — color cells"


# ── One deliverable, one file ───────────────────────────────────────────


def test_every_document_tool_states_the_same_path_rule():
    """Iterating on one deliverable must not mint a filename per attempt
    (the four-strays-for-one-document bug): every document-producing tool
    description carries the reuse-the-same-path rule."""
    from server.documents.tools import DOCUMENT_TOOLS

    producing = [t for t in DOCUMENT_TOOLS if t["name"] != "inspect_document"]
    assert len(producing) == 5
    for tool in producing:
        assert "ONE DELIVERABLE, ONE FILE" in tool["description"], tool["name"]


def test_office_script_advertises_in_place_revision():
    from server.documents.tools import OFFICE_SCRIPT_TOOL

    desc = OFFICE_SCRIPT_TOOL["description"]
    assert "MAY be the same file as source_path" in desc
    src = OFFICE_SCRIPT_TOOL["input_schema"]["properties"]["source_path"]
    assert "same file as the destination" in src["description"]


def test_write_outcome_points_the_next_revision_at_the_same_path(tmp_path, monkeypatch):
    """The write result is read at the exact moment the model picks its next
    attempt's path — it must say: same destination, replaces the file."""
    monkeypatch.setattr("server.workspace.get_workspace_path", lambda: str(tmp_path))
    outcome = _run(_do_office_script(_payload(tmp_path, _MAKE_DOCX)))
    assert outcome.ok, outcome.error
    assert "same destination" in outcome.output
    assert "unless the user asks for a separate copy" in outcome.output


def test_revising_in_place_replaces_rather_than_multiplies(tmp_path):
    """source_path == destination is the supported revision flow: the input
    is copied before the run, the output replaces the file, and the folder
    ends up holding exactly one file."""
    dest = str(tmp_path / "report.docx")
    build = (
        "from docx import Document\nd = Document()\nd.add_paragraph('v1')\nd.save(OUTPUT_PATH)\n"
    )
    outcome = _run(_do_office_script(_payload(tmp_path, build, "absolute", dest)))
    assert outcome.ok, outcome.error

    revise = (
        "from docx import Document\n"
        "d = Document(INPUT_PATH)\n"
        "d.add_paragraph('v2')\n"
        "d.save(OUTPUT_PATH)\n"
    )
    outcome = _run(
        _do_office_script(_payload(tmp_path, revise, "absolute", dest, source_path=dest))
    )
    assert outcome.ok, outcome.error

    from docx import Document

    assert [p.text for p in Document(dest).paragraphs] == ["v1", "v2"]
    assert [p.name for p in tmp_path.iterdir()] == ["report.docx"]
