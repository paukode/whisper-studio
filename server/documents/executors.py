"""@register_executor handlers for the document-creation tools.

The tools NEVER ask the user where to save and never force a workspace
connection. Destination resolution is the shared
server.workspace.state.resolve_write_destination (also used by
ws_create_file/save_file): destination_path from the user's own words wins;
otherwise a connected workspace's `path`; otherwise ~/Documents/<filename>.

The document is fully built at tool-call time. For a workspace destination
the bytes ride in a "create_document" approval payload (server/approval/
bootstrap.py's _do_create_document re-resolves the workspace root at click
time). For an absolute destination the finished bytes are STAGED into the
app sandbox first (server/workspace/paths.py's stage_file_bytes) and the
"save_to_path" approval simply moves the ready file into place — by the
time the user sees the card the document already exists, so approving can
never fail on content problems.

Either way, nothing lands outside the sandbox until a human clicks Yes.
There is no diff to preview (the model supplies HTML/slide-data/rows/
paragraphs, not literal file bytes), so the approval card shows a plain
"Create <path> (<n> bytes)" confirmation (preview="text") rather than a
diff.

office_script is the exception to that build-then-approve shape. It carries
model-written Python rather than data, so it follows run_python's contract
instead: the executor only packages the code, the approval card shows it,
and it runs (in server/documents/script.py's sandbox) after the click. Its
read-only companion inspect_document needs no approval at all.
"""

import base64
import io
import json
import logging
import os
from xml.sax.saxutils import escape as _xml_escape

from server.documents.docx_html import render_html
from server.documents.inspect import inspect_document
from server.documents.script import resolve_source_path, validate_script_request
from server.documents.standards import (
    DOCX_BODY_FONT,
    DOCX_BODY_SIZE_PT,
    apply_docx_standard,
    apply_pptx_standard,
    apply_xlsx_standard,
    fit_pptx_placeholders,
    pdf_standard_styles,
)
from server.executors import register_executor
from server.workspace.executors import _stage_and_request_save
from server.workspace.state import resolve_write_destination

log = logging.getLogger("whisper-studio")

# Generous but bounded: this is HTML *source*, not the rendered document —
# a few MB of markup is already a very large document. Guards against
# handing the parser a runaway payload.
_MAX_HTML_CONTENT_BYTES = 5_000_000


def _resolve_target(tool_name: str, tool_input: dict) -> tuple[str, str] | str:
    """Thin wrapper around the shared resolve_write_destination (server/
    workspace/state.py), adding the is-a-directory check specific to
    document tools — a docx/pptx/xlsx/pdf can never target an existing
    directory, unlike the generic resolver's callers."""
    resolved = resolve_write_destination(tool_name, tool_input)
    if isinstance(resolved, str):
        return resolved
    kind, target = resolved
    if kind == "workspace" and os.path.isdir(target):
        path = tool_input.get("path", "")
        return f"Error: {path} is an existing directory, not a file path."
    return resolved


def _ensure_ext(path: str, fmt: str) -> str:
    """Append the format's extension when missing — models sometimes pass a
    bare name like 'document', and the deliverable must open in the right
    app regardless."""
    return path if path.lower().endswith(f".{fmt}") else f"{path}.{fmt}"


def _approval_result(resolved: tuple[str, str], relative_path: str, data: bytes, fmt: str) -> str:
    """Package the FINISHED document bytes for approval — the executor never
    writes to the destination directly.

    `resolved` is _resolve_target's success tuple. For ("workspace", ...),
    packages a create_document action keyed by the workspace-relative path.
    For ("absolute", ...), stages the bytes in the app sandbox and packages
    the shared save_to_path move action (see module docstring)."""
    kind, target = resolved
    if kind == "absolute":
        target = _ensure_ext(target, fmt)
        return _stage_and_request_save(target, os.path.basename(target), data)
    payload = json.dumps(
        {
            "action": "create_document",
            "path": _ensure_ext(relative_path, fmt),
            "content_b64": base64.b64encode(data).decode("ascii"),
            "format": fmt,
            "size": len(data),
        }
    )
    return f"[WS_APPROVAL]{payload}"


@register_executor("create_docx", read_only=False, concurrent_safe=False)
def _exec_create_docx(tool_input, transcript, current_attachments):
    resolved = _resolve_target("create_docx", tool_input)
    if isinstance(resolved, str):
        return resolved
    path = tool_input.get("path", "")

    html_content = tool_input.get("html_content", "")
    if not html_content or not html_content.strip():
        return "Error: html_content is required and must be non-empty."
    encoded_len = len(html_content.encode("utf-8", errors="ignore"))
    if encoded_len > _MAX_HTML_CONTENT_BYTES:
        return (
            f"Error: html_content is too large ({encoded_len} bytes) — keep it under "
            f"{_MAX_HTML_CONTENT_BYTES} bytes."
        )

    from docx import Document

    try:
        doc = Document()
        apply_docx_standard(doc)
        render_html(
            doc,
            html_content,
            body_font=DOCX_BODY_FONT,
            body_size_pt=DOCX_BODY_SIZE_PT,
        )
        buf = io.BytesIO()
        doc.save(buf)
    except Exception as e:
        return f"Error building docx: {e}"

    return _approval_result(resolved, path, buf.getvalue(), "docx")


@register_executor("create_pptx", read_only=False, concurrent_safe=False)
def _exec_create_pptx(tool_input, transcript, current_attachments):
    resolved = _resolve_target("create_pptx", tool_input)
    if isinstance(resolved, str):
        return resolved
    path = tool_input.get("path", "")

    slides = tool_input.get("slides")
    if not isinstance(slides, list) or not slides:
        return "Error: slides must be a non-empty list of {title, bullets}."
    for i, slide_spec in enumerate(slides):
        if not isinstance(slide_spec, dict):
            return f"Error: slides[{i}] must be an object with 'title' and 'bullets'."
        if not isinstance(slide_spec.get("bullets") or [], list):
            return f"Error: slides[{i}].bullets must be a list of strings."

    from pptx import Presentation

    try:
        prs = Presentation()
        apply_pptx_standard(prs)
        layout = prs.slide_layouts[1]  # "Title and Content"
        for slide_spec in slides:
            title = str(slide_spec.get("title", "") or "")
            bullets = slide_spec.get("bullets") or []
            slide = prs.slides.add_slide(layout)
            fit_pptx_placeholders(prs, slide)
            slide.shapes.title.text = title
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            for j, bullet in enumerate(bullets):
                para = text_frame.paragraphs[0] if j == 0 else text_frame.add_paragraph()
                para.text = str(bullet)
        buf = io.BytesIO()
        prs.save(buf)
    except Exception as e:
        return f"Error building pptx: {e}"

    return _approval_result(resolved, path, buf.getvalue(), "pptx")


@register_executor("create_xlsx", read_only=False, concurrent_safe=False)
def _exec_create_xlsx(tool_input, transcript, current_attachments):
    resolved = _resolve_target("create_xlsx", tool_input)
    if isinstance(resolved, str):
        return resolved
    path = tool_input.get("path", "")

    sheets = tool_input.get("sheets")
    if not isinstance(sheets, list) or not sheets:
        return "Error: sheets must be a non-empty list of {name, rows}."
    for i, sheet_spec in enumerate(sheets):
        if not isinstance(sheet_spec, dict):
            return f"Error: sheets[{i}] must be an object with 'name' and 'rows'."
        if not isinstance(sheet_spec.get("rows") or [], list):
            return f"Error: sheets[{i}].rows must be a list of rows."

    from openpyxl import Workbook

    try:
        wb = Workbook()
        for i, sheet_spec in enumerate(sheets):
            name = str(sheet_spec.get("name") or f"Sheet{i + 1}")
            rows = sheet_spec.get("rows") or []
            worksheet = wb.active if i == 0 else wb.create_sheet()
            worksheet.title = name
            for row in rows:
                if not isinstance(row, list):
                    return f"Error: sheets[{i}].rows must be a list of lists."
                worksheet.append(row)
            apply_xlsx_standard(worksheet)
        buf = io.BytesIO()
        wb.save(buf)
    except Exception as e:
        return f"Error building xlsx: {e}"

    return _approval_result(resolved, path, buf.getvalue(), "xlsx")


@register_executor("inspect_document", read_only=True, concurrent_safe=True)
def _exec_inspect_document(tool_input, transcript, current_attachments):
    """Read-only structural map of an existing Office file. No approval:
    it opens the file, reports its skeleton, and changes nothing."""
    ok, resolved = resolve_source_path(tool_input.get("path", ""))
    if not ok:
        return resolved
    return inspect_document(resolved)


@register_executor("office_script", read_only=False, concurrent_safe=False)
def _exec_office_script(tool_input, transcript, current_attachments):
    """Package model-written python-docx/pptx/openpyxl code for approval.

    Nothing runs here — the code executes only after the user clicks Yes
    (server/documents/approval.py's _do_office_script). What this DOES do is
    resolve and validate the request up front, so the errors a model can fix
    on its own (missing code, wrong extension, unreadable source file) come
    back as tool results instead of wasting the user's click."""
    resolved = _resolve_target("office_script", tool_input)
    if isinstance(resolved, str):
        return resolved
    kind, target = resolved
    path = tool_input.get("path", "")

    # The extension that decides the format is the one on the FINAL
    # destination: with destination_path set, that is the resolved target,
    # not the (possibly bare) suggested `path`.
    named = path if kind == "workspace" else target
    out_format = os.path.splitext(named)[1].lower().lstrip(".")
    code = tool_input.get("code", "")
    source_path = str(tool_input.get("source_path") or "").strip()
    error = validate_script_request(code, out_format, source_path)
    if error:
        return error

    payload = json.dumps(
        {
            "action": "office_script",
            "path": path if kind == "workspace" else target,
            "dest_kind": kind,
            "format": out_format,
            "code": code,
            "source_path": source_path,
            "summary": str(tool_input.get("summary") or "").strip(),
        }
    )
    return f"[WS_APPROVAL]{payload}"


@register_executor("create_pdf", read_only=False, concurrent_safe=False)
def _exec_create_pdf(tool_input, transcript, current_attachments):
    resolved = _resolve_target("create_pdf", tool_input)
    if isinstance(resolved, str):
        return resolved
    path = tool_input.get("path", "")

    title = str(tool_input.get("title", "") or "")
    paragraphs = tool_input.get("paragraphs")
    if not isinstance(paragraphs, list):
        return "Error: paragraphs must be a list of strings."
    if not title and not paragraphs:
        return "Error: provide a title and/or at least one paragraph."

    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    def _flowable_text(raw: str) -> str:
        # paragraphs/title are plain text, not reportlab's mini-XML markup —
        # escape so a stray "<", ">", or "&" can't be mistaken for a tag
        # (which silently drops or mangles the surrounding text) and turn a
        # literal newline into a real line break within the paragraph.
        return _xml_escape(raw).replace("\n", "<br/>")

    try:
        styles = pdf_standard_styles()
        flowables = []
        if title:
            flowables.append(Paragraph(_flowable_text(title), styles["Title"]))
            flowables.append(Spacer(1, 12))
        for i, para in enumerate(paragraphs):
            flowables.append(Paragraph(_flowable_text(str(para)), styles["Normal"]))
            if i != len(paragraphs) - 1:
                flowables.append(Spacer(1, 12))

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4)
        doc.build(flowables)
    except Exception as e:
        return f"Error building pdf: {e}"

    return _approval_result(resolved, path, buf.getvalue(), "pdf")
