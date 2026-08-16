"""Structural read-out of an existing .docx / .xlsx / .pptx.

The `office_script` tool lets the model edit a real Office file, but it can
only address what it can see: "shade the header row of the second table"
needs a table index, a row count, and a column count, not prose. The
attachment/extraction pipeline (server/extract/) answers a different
question — it flattens a document to markdown for reading — and throws away
exactly the coordinates an editor needs.

So this module reports the skeleton: table indices and dimensions, sheet
dimensions and merged ranges, slide shape inventories with positions. Cell
*text* is included only as a shallow sample, enough to identify which table
or column is which without turning a structural map back into a full
extraction.

Output is plain text, addressed in the coordinate system of the library
that will do the editing (python-docx `tables[i].rows[r].cells[c]`,
openpyxl `ws["B2"]`, python-pptx `slides[i].shapes[j]`), so a coordinate
read here can be pasted straight into a script.
"""

import os

# Every list in the report is capped: a 5000-row sheet or a 200-slide deck
# must still produce a report the model can read in one go. The caps are per
# section, so a big workbook degrades to "first N of M" rather than to a
# truncated blob with the tail silently missing.
MAX_TABLES = 20
MAX_ROWS_PER_TABLE = 8
MAX_SHEETS = 20
MAX_SHEET_ROWS = 8
MAX_SLIDES = 40
MAX_SHAPES_PER_SLIDE = 20
MAX_CELL_CHARS = 40

EMU_PER_INCH = 914400

SUPPORTED_FORMATS = ("docx", "xlsx", "pptx")


def _cell(value) -> str:
    """One cell's text, flattened to a single line and clipped. Newlines
    inside a cell would break the one-row-per-line report layout."""
    text = "" if value is None else str(value)
    text = " ".join(text.split())
    if len(text) > MAX_CELL_CHARS:
        text = text[: MAX_CELL_CHARS - 1] + "…"
    return text


def _row_line(label: str, values: list) -> str:
    return f"    {label}: {[_cell(v) for v in values]}"


def _inches(emu) -> str:
    if emu is None:
        return "?"
    return f"{emu / EMU_PER_INCH:.2f}in"


def _more(shown: int, total: int, noun: str) -> list[str]:
    if total > shown:
        return [f"  ... {total - shown} more {noun} not shown"]
    return []


def _inspect_docx(path: str) -> str:
    from docx import Document

    doc = Document(path)
    lines = [f"DOCX: {os.path.basename(path)}"]

    section = doc.sections[0] if doc.sections else None
    if section is not None:
        lines.append(
            f"Page: {_inches(section.page_width)} x {_inches(section.page_height)}, "
            f"margins L/R {_inches(section.left_margin)}/{_inches(section.right_margin)}, "
            f"T/B {_inches(section.top_margin)}/{_inches(section.bottom_margin)}"
        )
    lines.append(f"Paragraphs: {len(doc.paragraphs)}   Tables: {len(doc.tables)}")
    if doc.inline_shapes:
        lines.append(f"Inline images: {len(doc.inline_shapes)}")

    # Heading outline: the fastest way for the model to locate the section it
    # was asked to change without reading the whole body back.
    headings = [
        (p.style.name, p.text.strip())
        for p in doc.paragraphs
        if p.style is not None and p.style.name.startswith("Heading") and p.text.strip()
    ]
    if headings:
        lines.append("")
        lines.append("Headings:")
        for style_name, text in headings[:MAX_TABLES]:
            lines.append(f"  {style_name}: {_cell(text)}")
        lines.extend(_more(min(len(headings), MAX_TABLES), len(headings), "headings"))

    for i, table in enumerate(doc.tables[:MAX_TABLES]):
        rows = table.rows
        n_cols = len(table.columns)
        style_name = table.style.name if table.style is not None else "(none)"
        lines.append("")
        lines.append(f"Table {i} — {len(rows)} rows x {n_cols} cols, style '{style_name}'")
        lines.append(f"  address cells as: doc.tables[{i}].rows[r].cells[c]  (r,c are 0-based)")
        for r, row in enumerate(rows[:MAX_SHEET_ROWS]):
            lines.append(_row_line(f"row {r}", [c.text for c in row.cells]))
        lines.extend(_more(min(len(rows), MAX_SHEET_ROWS), len(rows), "rows"))
    lines.extend(_more(min(len(doc.tables), MAX_TABLES), len(doc.tables), "tables"))
    return "\n".join(lines)


def _inspect_xlsx(path: str) -> str:
    from openpyxl import load_workbook

    # data_only=False keeps formulas visible as formulas: an editor needs to
    # know a cell holds "=SUM(B2:B9)" before overwriting it with a number.
    wb = load_workbook(path, data_only=False)
    lines = [f"XLSX: {os.path.basename(path)}", f"Sheets: {len(wb.sheetnames)}"]

    for name in wb.sheetnames[:MAX_SHEETS]:
        ws = wb[name]
        lines.append("")
        lines.append(
            f"Sheet '{name}' — {ws.max_row} rows x {ws.max_column} cols "
            f"(used range {ws.dimensions})"
        )
        lines.append(f"  address cells as: wb['{name}']['B2']  or  .cell(row=2, column=2)")
        if ws.freeze_panes:
            lines.append(f"  frozen panes at {ws.freeze_panes}")
        merged = [str(r) for r in ws.merged_cells.ranges]
        if merged:
            shown = merged[:MAX_SHEET_ROWS]
            lines.append(f"  merged ranges: {', '.join(shown)}")
            lines.extend(_more(len(shown), len(merged), "merged ranges"))
        for r, row in enumerate(
            ws.iter_rows(max_row=min(ws.max_row, MAX_SHEET_ROWS), values_only=True), start=1
        ):
            lines.append(_row_line(f"row {r}", list(row)))
        lines.extend(_more(min(ws.max_row, MAX_SHEET_ROWS), ws.max_row, "rows"))
    lines.extend(_more(min(len(wb.sheetnames), MAX_SHEETS), len(wb.sheetnames), "sheets"))
    return "\n".join(lines)


def _shape_line(index: int, shape) -> str:
    kind = "shape"
    try:
        if shape.is_placeholder:
            kind = f"placeholder {shape.placeholder_format.type}"
        elif shape.shape_type is not None:
            kind = str(shape.shape_type)
    except (AttributeError, ValueError):
        # python-pptx raises on placeholders/shapes whose type it cannot map;
        # the geometry below is still useful, so fall back to the generic label.
        pass
    text = ""
    if getattr(shape, "has_text_frame", False):
        text = f' "{_cell(shape.text_frame.text)}"'
    if getattr(shape, "has_table", False):
        table = shape.table
        text = f" table {len(table.rows)}x{len(table.columns)}"
    geometry = (
        f"at ({_inches(shape.left)}, {_inches(shape.top)}) "
        f"{_inches(shape.width)} x {_inches(shape.height)}"
    )
    return f"  [{index}] {kind} '{shape.name}'{text} {geometry}"


def _inspect_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    lines = [
        f"PPTX: {os.path.basename(path)}",
        f"Slide size: {_inches(prs.slide_width)} x {_inches(prs.slide_height)}",
        f"Slides: {len(prs.slides)}",
        "Available layouts: "
        + ", ".join(f"[{i}] {layout.name}" for i, layout in enumerate(prs.slide_layouts)),
    ]

    for i, slide in enumerate(list(prs.slides)[:MAX_SLIDES]):
        shapes = list(slide.shapes)
        lines.append("")
        lines.append(
            f"Slide {i} — layout '{slide.slide_layout.name}', {len(shapes)} shapes"
            f"   (address as prs.slides[{i}].shapes[j])"
        )
        for j, shape in enumerate(shapes[:MAX_SHAPES_PER_SLIDE]):
            lines.append(_shape_line(j, shape))
        lines.extend(_more(min(len(shapes), MAX_SHAPES_PER_SLIDE), len(shapes), "shapes"))
    lines.extend(_more(min(len(prs.slides), MAX_SLIDES), len(prs.slides), "slides"))
    return "\n".join(lines)


_INSPECTORS = {
    "docx": _inspect_docx,
    "xlsx": _inspect_xlsx,
    "pptx": _inspect_pptx,
}


def format_of(path: str) -> str:
    """The supported format for a path, or "" when the extension isn't one.
    Shared with the script runner so both tools agree on what "an Office
    file" means (.doc/.xls/.ppt are the pre-2007 binary formats, which none
    of these libraries can open)."""
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    return ext if ext in SUPPORTED_FORMATS else ""


def inspect_document(path: str) -> str:
    """Structural report for one Office file, or an "Error: ..." string.

    Errors are returned rather than raised: every caller is a tool executor
    whose contract is to hand the model a string it can act on."""
    fmt = format_of(path)
    if not fmt:
        return (
            f"Error: {os.path.basename(path)} is not a .docx/.xlsx/.pptx file. "
            "The pre-2007 .doc/.xls/.ppt formats are not supported — open the file "
            "in Word/Excel/PowerPoint and re-save it in the modern format first."
        )
    if not os.path.isfile(path):
        return f"Error: no such file: {path}"
    try:
        return _INSPECTORS[fmt](path)
    except Exception as e:
        return f"Error reading {os.path.basename(path)}: {e}"
